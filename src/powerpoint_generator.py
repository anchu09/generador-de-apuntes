"""
Módulo para generar PowerPoint con diapositivas, transcripciones y resúmenes.
"""

import re
from pathlib import Path
from typing import List, Tuple
from PIL import Image, ImageDraw, ImageFont
from .config_loader import PowerPointConfig
from .video_processor import SlideSegment


class PowerPointGenerator:
    """
    Genera presentaciones PowerPoint con diapositivas originales,
    transcripciones y resúmenes.
    """

    def __init__(self, config: PowerPointConfig):
        """
        Inicializa el generador de PowerPoint.

        Args:
            config: Configuración de generación de PowerPoint
        """
        self.config = config

    def generate_presentation(
        self, segments: List[SlideSegment], output_path: str
    ) -> str:
        """
        Genera una presentación PowerPoint o PDF a partir de los segmentos.

        Args:
            segments: Lista de segmentos con diapositivas, transcripciones y resúmenes
            output_path: Ruta donde guardar la presentación

        Returns:
            Ruta al archivo generado
        """
        output_format = self.config.output_format.lower()

        if output_format == "pdf":
            return self._generate_pdf(segments, output_path)
        elif output_format == "pptx":
            return self._generate_pptx(segments, output_path)
        else:
            raise ValueError(
                f"Formato no soportado: {output_format}. Use 'pdf' o 'pptx'"
            )

    def _generate_pdf(self, segments: List[SlideSegment], output_path: str) -> str:
        """Genera un PDF a partir de los segmentos."""
        # Calcular tamaño de diapositiva (ancho x 3)
        # Para PDF no hay límite de tamaño como en PowerPoint
        slide_width = self.config.slide_width * 3
        slide_height = self.config.slide_height

        # Calcular ancho de cada columna (tercio del ancho total)
        column_width = slide_width // 3

        # Asegurar que el directorio de salida existe
        output_path_obj = Path(output_path)
        output_path_obj.parent.mkdir(parents=True, exist_ok=True)

        # Lista para almacenar todas las imágenes compuestas
        composite_images = []

        for segment in segments:
            # Cargar imagen de la diapositiva original y escalarla al tamaño de columna
            if not Path(segment.frame_path).exists():
                print(
                    f"  Advertencia: No se encuentra {segment.frame_path}, saltando..."
                )
                continue

            slide_img = Image.open(segment.frame_path).convert("RGB")
            slide_img = self._resize_to_fit(slide_img, column_width, slide_height)
            offset_x = (column_width - slide_img.width) // 2
            offset_y = (slide_height - slide_img.height) // 2

            # Crear imagen compuesta
            composite_width = slide_width
            composite_height = slide_height
            composite = Image.new("RGB", (composite_width, composite_height), "white")

            # Colocar diapositiva original en el centro
            center_x = column_width + offset_x
            composite.paste(slide_img, (center_x, offset_y))

            # Agregar texto de transcripción a la izquierda (si existe)
            if segment.transcription:
                # Formatear transcripción (párrafos, puntos y aparte)
                formatted_text = self._format_transcription(segment.transcription)

                # Calcular tamaño de fuente óptimo si está habilitado
                font_size = self.config.font_size_transcription
                if self.config.auto_font_size:
                    font_size = self._calculate_optimal_font_size(
                        formatted_text, column_width, slide_height
                    )

                self._add_text_to_image(
                    composite,
                    formatted_text,
                    x=0,
                    y=0,
                    width=column_width,
                    height=slide_height,
                    font_size=font_size,
                    align="left",
                )
            else:
                # Si no hay transcripción, dejar espacio en blanco
                pass

            # Agregar resumen a la derecha (si existe)
            if segment.summary:
                # Formatear resumen (párrafos, puntos y aparte)
                formatted_summary = self._format_transcription(segment.summary)

                # Calcular tamaño de fuente óptimo si está habilitado
                font_size = self.config.font_size_summary
                if self.config.auto_font_size:
                    font_size = self._calculate_optimal_font_size(
                        formatted_summary, column_width, slide_height
                    )

                self._add_text_to_image(
                    composite,
                    formatted_summary,
                    x=column_width * 2,  # Columna derecha
                    y=0,
                    width=column_width,
                    height=slide_height,
                    font_size=font_size,
                    align="left",
                )

            composite_images.append(composite)

        # Guardar todas las imágenes como PDF
        if composite_images:
            # Convertir a RGB si es necesario (PDF requiere RGB)
            rgb_images = []
            for img in composite_images:
                if img.mode != "RGB":
                    rgb_images.append(img.convert("RGB"))
                else:
                    rgb_images.append(img)

            # Guardar como PDF (primera imagen como base, resto como páginas adicionales)
            rgb_images[0].save(
                str(output_path),
                "PDF",
                resolution=100.0,
                save_all=True,
                append_images=rgb_images[1:] if len(rgb_images) > 1 else [],
            )
            print(f"  ✓ PDF generado con {len(composite_images)} páginas")
        else:
            print(f"  ⚠ No hay segmentos para generar PDF")

        return str(output_path)

    def _generate_pptx(self, segments: List[SlideSegment], output_path: str) -> str:
        """Genera un PowerPoint a partir de los segmentos."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise ImportError(
                "python-pptx no está instalado. Instala con: pip install python-pptx"
            )

        # Crear presentación
        prs = Presentation()

        # Configurar tamaño de diapositiva (ancho x 3)
        # PowerPoint tiene un límite máximo de 56 pulgadas (5376 píxeles a 96 DPI)
        max_width_pixels = 5376  # 56 pulgadas * 96 DPI
        desired_width = self.config.slide_width * 3

        if desired_width > max_width_pixels:
            # Escalar proporcionalmente para que quepa
            scale_factor = max_width_pixels / desired_width
            slide_width = max_width_pixels
            slide_height = int(self.config.slide_height * scale_factor)
        else:
            slide_width = desired_width
            slide_height = self.config.slide_height

        prs.slide_width = Inches(slide_width / 96)  # Convertir píxeles a pulgadas
        prs.slide_height = Inches(slide_height / 96)

        # Calcular ancho de cada columna (tercio del ancho total)
        column_width = slide_width // 3

        for segment in segments:
            if not segment.transcription:
                continue

            # Crear nueva diapositiva
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout en blanco

            # Cargar imagen de la diapositiva original y escalarla al tamaño de columna
            slide_img = Image.open(segment.frame_path).convert("RGB")
            slide_img = self._resize_to_fit(slide_img, column_width, slide_height)
            offset_x = (column_width - slide_img.width) // 2
            offset_y = (slide_height - slide_img.height) // 2

            # Crear imagen compuesta
            composite_width = slide_width
            composite_height = slide_height
            composite = Image.new("RGB", (composite_width, composite_height), "white")

            # Colocar diapositiva original en el centro
            center_x = column_width + offset_x
            composite.paste(slide_img, (center_x, offset_y))

            # Agregar texto de transcripción a la izquierda
            # Formatear transcripción (párrafos, puntos y aparte)
            formatted_text = self._format_transcription(segment.transcription)

            # Calcular tamaño de fuente óptimo si está habilitado
            font_size = self.config.font_size_transcription
            if self.config.auto_font_size:
                font_size = self._calculate_optimal_font_size(
                    formatted_text, column_width, slide_height
                )

            self._add_text_to_image(
                composite,
                formatted_text,
                x=0,
                y=0,
                width=column_width,
                height=slide_height,
                font_size=font_size,
                align="left",
            )

            # Agregar resumen a la derecha (si existe)
            if segment.summary:
                # Formatear resumen (párrafos, puntos y aparte)
                formatted_summary = self._format_transcription(segment.summary)

                # Calcular tamaño de fuente óptimo si está habilitado
                font_size = self.config.font_size_summary
                if self.config.auto_font_size:
                    font_size = self._calculate_optimal_font_size(
                        formatted_summary, column_width, slide_height
                    )

                self._add_text_to_image(
                    composite,
                    formatted_summary,
                    x=column_width * 2,  # Columna derecha
                    y=0,
                    width=column_width,
                    height=slide_height,
                    font_size=font_size,
                    align="left",
                )

            # Guardar imagen temporal
            temp_img_path = (
                Path(segment.frame_path).parent
                / f"temp_composite_{segment.slide_number}.png"
            )
            composite.save(temp_img_path)

            # Agregar imagen a la diapositiva
            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(
                str(temp_img_path),
                left,
                top,
                width=Inches(slide_width / 96),
                height=Inches(slide_height / 96),
            )

            # Eliminar imagen temporal
            temp_img_path.unlink()

        # Guardar presentación
        prs.save(output_path)
        return output_path

    def _format_transcription(self, text: str) -> str:
        """
        Formatea la transcripción para mejorar la legibilidad.
        Detecta pausas naturales, crea párrafos y añade saltos de línea.

        Args:
            text: Texto sin formatear

        Returns:
            Texto formateado con párrafos y saltos de línea
        """
        if not text:
            return text

        text = self._normalize_markdown(text)

        # Normalizar espacios múltiples
        text = re.sub(r"\s+", " ", text.strip())

        # Detectar pausas naturales: puntos, signos de interrogación/exclamación
        # Seguidos de espacio y mayúscula (inicio de nueva frase)
        # Patrón: punto/interrogación/exclamación + espacio + mayúscula
        text = re.sub(r"([.!?])\s+([A-ZÁÉÍÓÚÑ])", r"\1\n\n\2", text)

        # También detectar comas seguidas de pausas largas (más de 3 palabras después)
        # Esto es más conservador, solo para casos obvios
        # Dividir por comas y detectar frases muy largas
        sentences = re.split(r"([.!?])", text)
        formatted_sentences = []

        for i in range(0, len(sentences), 2):
            if i < len(sentences):
                sentence = sentences[i]
                punctuation = sentences[i + 1] if i + 1 < len(sentences) else ""

                # Si la frase es muy larga (más de 100 caracteres), buscar comas para dividir
                if len(sentence) > 100:
                    # Dividir por comas pero mantener la estructura
                    parts = sentence.split(",")
                    if len(parts) > 1:
                        # Reunir partes pero añadir saltos en comas estratégicas
                        sentence = ",\n".join(parts[:-1]) + "," + parts[-1]

                formatted_sentences.append(sentence + punctuation)

        text = "".join(formatted_sentences)

        # Limpiar saltos de línea múltiples (máximo 2 seguidos)
        text = re.sub(r"\n{3,}", "\n\n", text)

        # Eliminar espacios al inicio de líneas
        lines = text.split("\n")
        lines = [line.lstrip() for line in lines]
        text = "\n".join(lines)

        return text.strip()

    def _normalize_markdown(self, text: str) -> str:
        """
        Limpia formato markdown simple (negritas/listas) para mejorar el renderizado.
        """
        if not text:
            return text

        cleaned = text

        cleaned = re.sub(r"\*\*(.*?)\*\*", r"\1", cleaned)
        cleaned = re.sub(r"__(.*?)__", r"\1", cleaned)
        cleaned = re.sub(r"\*(.*?)\*", r"\1", cleaned)

        cleaned = re.sub(r"\s*-\s+", "\n• ", cleaned)
        cleaned = re.sub(r"\s*\*\s+", "\n• ", cleaned)

        cleaned = re.sub(r"\s*(\d+\.)\s+", r"\n\1 ", cleaned)

        cleaned = re.sub(r"\n{2,}", "\n", cleaned)

        return cleaned

    def _calculate_optimal_font_size(
        self, text: str, available_width: int, available_height: int
    ) -> int:
        """
        Calcula el tamaño de fuente óptimo para que el texto quepa en el espacio disponible.

        Args:
            text: Texto a renderizar
            available_width: Ancho disponible en píxeles
            available_height: Alto disponible en píxeles

        Returns:
            Tamaño de fuente óptimo (entre min_font_size y max_font_size)
        """
        # Crear una imagen temporal más grande para medir texto largo
        temp_img = Image.new("RGB", (available_width, available_height), "white")
        temp_draw = ImageDraw.Draw(temp_img)

        # Estimar líneas basándose en el ancho disponible
        # Usar búsqueda binaria para encontrar el tamaño óptimo
        min_size = self.config.min_font_size
        max_size = self.config.max_font_size
        # Empezar con un tamaño más grande para textos cortos
        word_count = len(text.split())
        if word_count < 50:
            initial_size = max_size
        elif word_count < 150:
            initial_size = (min_size + max_size) // 2
        else:
            initial_size = self.config.font_size_transcription
        optimal_size = initial_size

        # Búsqueda binaria para encontrar el tamaño óptimo
        left, right = min_size, max_size

        while left <= right:
            test_size = (left + right) // 2
            try:
                test_font = ImageFont.truetype(
                    "/System/Library/Fonts/Helvetica.ttc", test_size
                )
            except:
                try:
                    test_font = ImageFont.truetype("arial.ttf", test_size)
                except:
                    test_font = ImageFont.load_default()

            # Calcular cuántas líneas necesitaría el texto
            # Respetar saltos de línea existentes (párrafos)
            paragraphs = text.split("\n\n")
            lines = []
            effective_width = available_width - (2 * self.config.margin)

            for para in paragraphs:
                para_words = para.split()
                current_line = []

                for word in para_words:
                    test_line = (
                        " ".join(current_line + [word]) if current_line else word
                    )
                    bbox = temp_draw.textbbox((0, 0), test_line, font=test_font)
                    text_width = bbox[2] - bbox[0]

                    if text_width <= effective_width:
                        current_line.append(word)
                    else:
                        if current_line:
                            lines.append(" ".join(current_line))
                        current_line = [word]

                if current_line:
                    lines.append(" ".join(current_line))

                # Añadir línea vacía entre párrafos (excepto después del último)
                if para != paragraphs[-1]:
                    lines.append("")

            # Calcular altura total necesaria
            line_height = int(test_size * self.config.line_spacing)
            # Contar párrafos reales (no líneas vacías)
            num_paragraphs = text.count("\n\n") + 1
            # Solo contar líneas vacías como espaciado, no como líneas que ocupan espacio
            num_empty_lines = sum(1 for line in lines if not line.strip())
            num_text_lines = len(lines) - num_empty_lines
            paragraph_spacing = self.config.paragraph_spacing * max(
                0, num_paragraphs - 1
            )

            # Calcular altura: líneas de texto + espaciado entre párrafos + márgenes
            total_height = (
                num_text_lines * line_height
                + num_empty_lines
                * (line_height // 2)  # Líneas vacías ocupan menos espacio
                + paragraph_spacing
                + (2 * self.config.margin)
            )

            # Si el texto cabe con este tamaño, actualizar optimal_size y buscar uno mayor
            # Usar un margen de seguridad del 99% para usar casi todo el espacio disponible
            if total_height <= available_height * 0.99:
                optimal_size = test_size
                left = test_size + 1  # Intentar un tamaño mayor
            else:
                # Si no cabe, probar un tamaño menor
                right = test_size - 1

        # Asegurar tamaños mínimos según cantidad de texto
        # Cada diapositiva tendrá su tamaño optimizado según su contenido
        if word_count < 30:
            # Para textos muy cortos, usar un tamaño grande (al menos 80% del máximo)
            optimal_size = max(optimal_size, int(self.config.max_font_size * 0.8))
        elif word_count < 80:
            # Para textos cortos, usar al menos 60% del máximo
            optimal_size = max(optimal_size, int(self.config.max_font_size * 0.6))
        elif word_count < 200:
            # Para textos medianos, usar al menos 40% del máximo
            optimal_size = max(optimal_size, int(self.config.max_font_size * 0.4))

        return max(min_size, min(max_size, optimal_size))

    def _add_text_to_image(
        self,
        image: Image.Image,
        text: str,
        x: int,
        y: int,
        width: int,
        height: int,
        font_size: int,
        align: str = "left",
    ):
        """
        Agrega texto a una imagen.

        Args:
            image: Imagen PIL donde agregar el texto
            text: Texto a agregar
            x: Posición X inicial
            y: Posición Y inicial
            width: Ancho del área de texto
            height: Alto del área de texto
            font_size: Tamaño de fuente
            align: Alineación del texto ('left', 'right', 'center')
        """
        draw = ImageDraw.Draw(image)

        # Intentar cargar fuente, usar default si no está disponible
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", font_size)
        except:
            try:
                font = ImageFont.truetype("arial.ttf", font_size)
            except:
                font = ImageFont.load_default()

        # Dividir texto en líneas, respetando saltos de línea existentes
        # Primero dividir por saltos de línea (párrafos)
        paragraphs = text.split("\n\n")
        lines = []

        for para in paragraphs:
            # Para cada párrafo, dividir en líneas que quepan en el ancho
            words = para.split()
            current_line = []
            effective_width = width - (2 * self.config.margin)

            for word in words:
                test_line = " ".join(current_line + [word]) if current_line else word
                bbox = draw.textbbox((0, 0), test_line, font=font)
                text_width = bbox[2] - bbox[0]

                if text_width <= effective_width:
                    current_line.append(word)
                else:
                    if current_line:
                        lines.append(" ".join(current_line))
                    current_line = [word]

            if current_line:
                lines.append(" ".join(current_line))

            # Añadir línea vacía entre párrafos (excepto después del último)
            if para != paragraphs[-1]:
                lines.append("")

        # Dibujar líneas de texto con espaciado entre párrafos
        y_offset = y + self.config.margin
        line_height = int(font_size * self.config.line_spacing)
        paragraph_spacing = self.config.paragraph_spacing

        for i, line in enumerate(lines):
            # Verificar si hay espacio disponible
            if y_offset + line_height > y + height - self.config.margin:
                break

            # Detectar si es inicio de párrafo (línea vacía)
            is_paragraph_start = i > 0 and not line.strip()

            # Añadir espaciado extra si es inicio de párrafo
            if is_paragraph_start:
                y_offset += paragraph_spacing
                continue  # Saltar la línea vacía

            # Calcular posición X según alineación
            bbox = draw.textbbox((0, 0), line, font=font)
            text_width = bbox[2] - bbox[0]

            if align == "center":
                text_x = x + (width - text_width) // 2
            elif align == "right":
                text_x = x + width - text_width - self.config.margin
            else:  # left
                text_x = x + self.config.margin

            draw.text((text_x, y_offset), line, fill="black", font=font)
            y_offset += line_height

    def _resize_to_fit(
        self, image: Image.Image, max_width: int, max_height: int
    ) -> Image.Image:
        """
        Escala una imagen manteniendo su relación de aspecto para que encaje en un cuadro.
        """
        if not image.width or not image.height:
            return image

        scale = min(max_width / image.width, max_height / image.height)
        scale = max(scale, 1e-6)

        new_width = max(1, int(image.width * scale))
        new_height = max(1, int(image.height * scale))

        if (new_width, new_height) == image.size:
            return image

        resampling = (
            Image.Resampling.LANCZOS
            if hasattr(Image, "Resampling")
            else Image.LANCZOS
        )
        return image.resize((new_width, new_height), resampling)
